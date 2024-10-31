from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def set_font(paragraph, font_size, text_color):
    """ Helper function to set font size and color for a paragraph """
    run = paragraph.add_run()
    font = run.font
    font.size = Pt(font_size)
    
    # Set text color
    if text_color.lower() == "black":
        font.color.rgb = RGBColor(0, 0, 0)
    elif text_color.lower() == "dark blue":
        font.color.rgb = RGBColor(0, 0, 139)
    
    return run

def add_content_with_hierarchy(slide, slide_data):
    """ Helper function to add titles, subtitles, and content with hierarchy """
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()

    for section in slide_data["Formatted Content"]:
        # Add Subtitle (as second-level heading)
        subtitle_paragraph = text_frame.add_paragraph()
        run_subtitle = set_font(subtitle_paragraph, 20, section["Text Color"])
        run_subtitle.text = section.get("Subtitle", "")
        subtitle_paragraph.level = 0  # Subtitle level

        # Add content (as third-level bullet points)
        for content_line in section["Content"]:
            content_paragraph = text_frame.add_paragraph()
            content_paragraph.level = 1  # Content level (one level below subtitle)
            run_content = set_font(content_paragraph, int(section["Font Size"].replace("pt", "")), section["Text Color"])
            run_content.text = content_line

def prevent_overflow(slide):
    """ Adjust font sizes if content is too large for the slide by comparing text box size """
    text_frame = slide.shapes.placeholders[1].text_frame
    max_height = slide.shapes.placeholders[1].height  # Get the max height allowed by the text box

    def get_text_height(text_frame):
        """ Dummy function to estimate text height """
        # This is a simplified version and does not precisely compute height, but works as an approximation
        num_lines = sum([len(p.text.splitlines()) for p in text_frame.paragraphs])
        return Pt(20 * num_lines)  # Assumes 20pt height per line (adjust if necessary)

    text_height = get_text_height(text_frame)

    # Reduce font size while the text height exceeds the available space
    while text_height > max_height:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                current_size = run.font.size
                if current_size and current_size.pt > 10:  # Don't shrink too small
                    new_size = Pt(current_size.pt - 1)
                    run.font.size = new_size
        # Recalculate text height after adjusting font size
        text_height = get_text_height(text_frame)

def create_citation_slide(presentation, citations):
    """ Create a separate slide for all citations """
    slide_layout = presentation.slide_layouts[1]  # Title and content layout
    slide = presentation.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Citations"
    title_format = title.text_frame.paragraphs[0]
    title_format.font.size = Pt(36)
    title_format.font.color.rgb = RGBColor(0, 0, 139)  # Dark blue

    # Add citations content
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()

    for i, citation in enumerate(citations):
        p = text_frame.add_paragraph()
        run = set_font(p, 12, "Black")
        run.text = f"[{i+1}] {citation}"

def create_presentation(slides_data):
    # Create a new PowerPoint presentation object
    presentation = Presentation()
    citations = []

    for slide_data in slides_data:
        # Add a new slide with a title and content layout
        slide_layout = presentation.slide_layouts[1]  # Using layout with title and content
        slide = presentation.slides.add_slide(slide_layout)

        # Set the slide title and style it according to specifications
        title = slide.shapes.title
        title.text = slide_data["Slide Title"]
        
        # Set the title font size and color
        title_format = title.text_frame.paragraphs[0]
        title_format.font.size = Pt(int(slide_data["Font Size"].replace("pt", "")))
        if slide_data["Text Color"].lower() == "dark blue":
            title_format.font.color.rgb = RGBColor(0, 0, 139)  # Dark blue color
        
        # Add the formatted content with hierarchy
        add_content_with_hierarchy(slide, slide_data)

        # Collect citations for the citation slide
        if "Citations" in slide_data and slide_data["Citations"]:
            for citation in slide_data["Citations"]:
                citations.append(citation)

        # Prevent overflow by adjusting the font size if necessary
        prevent_overflow(slide)

    # Add a citation slide at the end if citations exist
    if citations:
        create_citation_slide(presentation, citations)

    # Save the presentation to a file
    presentation.save('Neuron_Solutions_Presentation.pptx')
    print("Presentation created successfully!")




import json

# Load the JSON data from the uploaded file
json_file_path = 'revised_slides.json'
with open(json_file_path, 'r') as file:
    slides_data = json.load(file)

# Create the presentation
create_presentation(slides_data)
